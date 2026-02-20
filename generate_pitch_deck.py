#!/usr/bin/env python3
"""
Perff AI — Investor Pitch Deck Generator
Generates a fully editable .pptx PowerPoint file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Brand Colors ──────────────────────────────────────────
DARK_BG     = RGBColor(0x0B, 0x15, 0x19)
DARK2_BG    = RGBColor(0x0E, 0x1E, 0x24)
GREEN       = RGBColor(0x63, 0xC2, 0x9D)
GREEN_DARK  = RGBColor(0x4A, 0x9E, 0x7D)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT  = RGBColor(0xD8, 0xE6, 0xEC)
TEXT_DIM    = RGBColor(0x58, 0x7E, 0x8C)
RED         = RGBColor(0xE8, 0x51, 0x5E)
ORANGE      = RGBColor(0xF0, 0xA5, 0x4B)
BLUE        = RGBColor(0x4B, 0xA3, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ── Helpers ───────────────────────────────────────────────

def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_paragraph(text_frame, text="", font_size=14, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(4), space_after=Pt(4),
                  font_name="Calibri"):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_slide_number(slide, num, total=15):
    add_textbox(slide, Inches(11.8), Inches(0.3), Inches(1.2), Inches(0.3),
                f"{num:02d} / {total}", font_size=10, color=TEXT_DIM,
                alignment=PP_ALIGN.RIGHT, font_name="Consolas")


def add_kicker(slide, text, left=Inches(0.8), top=Inches(0.8)):
    add_textbox(slide, left, top, Inches(6), Inches(0.3),
                text.upper(), font_size=10, color=GREEN, bold=True,
                font_name="Consolas")


def add_title(slide, line1, line2="", top=Inches(1.2), left=Inches(0.8),
              width=Inches(11), accent_color=GREEN):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(1.4))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    # First line in white
    run1 = p.add_run()
    run1.text = line1
    run1.font.size = Pt(40)
    run1.font.color.rgb = WHITE
    run1.font.bold = True
    run1.font.name = "Calibri"

    if line2:
        run_br = p.add_run()
        run_br.text = "\n"
        run2 = p.add_run()
        run2.text = line2
        run2.font.size = Pt(40)
        run2.font.color.rgb = accent_color
        run2.font.bold = True
        run2.font.name = "Calibri"

    return txBox


def add_title_centered(slide, line1, line2="", top=Inches(1.2)):
    txBox = add_title(slide, line1, line2, top, left=Inches(0.8), width=Inches(11.7))
    tf = txBox.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return txBox


def add_subtitle(slide, text, top=Inches(2.8), left=Inches(0.8), width=Inches(8)):
    return add_textbox(slide, left, top, width, Inches(1),
                       text, font_size=16, color=TEXT_LIGHT)


def add_stat_box(slide, left, top, value, label, badge="", width=Inches(1.8)):
    # Value
    add_textbox(slide, left, top, width, Inches(0.5),
                value, font_size=28, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")
    # Label
    add_textbox(slide, left, top + Inches(0.45), width, Inches(0.3),
                label, font_size=9, color=TEXT_DIM,
                alignment=PP_ALIGN.CENTER)
    if badge:
        add_textbox(slide, left, top + Inches(0.7), width, Inches(0.2),
                    badge, font_size=7, color=GREEN,
                    alignment=PP_ALIGN.CENTER, font_name="Consolas")


def add_card(slide, left, top, width, height, icon, title, body, stat=""):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    # Green top accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  left, top, width, Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()

    y = top + Inches(0.2)
    # Icon
    add_textbox(slide, left + Inches(0.15), y, Inches(0.5), Inches(0.35),
                icon, font_size=20, color=WHITE)
    y += Inches(0.35)
    # Title
    add_textbox(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.3),
                title, font_size=13, color=WHITE, bold=True)
    y += Inches(0.3)
    # Body
    add_textbox(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.6),
                body, font_size=10, color=TEXT_DIM)
    if stat:
        y += Inches(0.55)
        add_textbox(slide, left + Inches(0.15), y, width - Inches(0.3), Inches(0.3),
                    stat, font_size=18, color=GREEN, bold=True, font_name="Consolas")


def add_confidential(slide):
    add_textbox(slide, Inches(4.5), Inches(7.1), Inches(4.5), Inches(0.3),
                "CONFIDENTIAL — Perff AI © 2026", font_size=8,
                color=TEXT_DIM, alignment=PP_ALIGN.CENTER, font_name="Consolas")


def add_table(slide, left, top, width, col_widths, headers, rows):
    """Add a native PowerPoint table."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, Inches(0.3 * n_rows))
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(9)
            p.font.color.rgb = GREEN
            p.font.bold = True
            p.font.name = "Consolas"
            p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BG

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_LIGHT
                p.font.name = "Calibri"
            # Highlight first row (Perff AI)
            if r_idx == 0 and "Perff" in str(row[0]):
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = GREEN
                    p.font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK2_BG if r_idx % 2 == 0 else DARK_BG

    return table_shape


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 01 — HOOK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide)
add_slide_number(slide, 1)
add_kicker(slide, "The Opportunity", left=Inches(1.5), top=Inches(1.0))
add_title_centered(slide, "MENA Spends Billions on Ads.",
                   "Almost None of It Is Performance-Based.", top=Inches(1.4))

add_textbox(slide, Inches(2), Inches(3.0), Inches(9), Inches(0.8),
            "The region's e-commerce market grows 15% annually — faster than any major market — "
            "yet captures barely 2% of the $18B global affiliate marketing industry. That gap is our opportunity.",
            font_size=15, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Stats
stats = [("15%", "MENA E-com Growth", "FASTEST GLOBALLY"),
         ("2%", "Affiliate Capture", "VS 15% IN WEST"),
         ("$1.5B", "MENA SAM", "GROWING 25%+ YoY"),
         ("335×", "CAC Premium", "VS AFFILIATE CHANNEL")]
x_start = Inches(1.5)
for i, (val, label, badge) in enumerate(stats):
    add_stat_box(slide, x_start + Inches(2.7 * i), Inches(4.2), val, label, badge)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 02 — WHAT WE DO
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 2)
add_kicker(slide, "What We Do")
add_title(slide, 'The Shopify of Affiliate Marketing',
          'for the Middle East.')

add_subtitle(slide,
    "Perff AI is an AI-powered platform connecting brands with affiliate publishers — "
    "automating campaign matching, fraud detection, and performance optimization so brands pay only for real results.")

# Flow diagram as text
flow_items = ["🏢 Brand → Sets budget & goals",
              "🧠 Perff AI → AI matches, optimizes, protects",
              "📱 Publisher → Promotes to audience",
              "👤 Customer → Converts (buys/signs up)",
              "💰 Everyone Paid → Performance-based. Automatic."]
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.5), Inches(0.6))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "  →  ".join([f.split(" → ")[0] for f in flow_items])
p.font.size = Pt(14)
p.font.color.rgb = GREEN
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Cards row
cards = [("🎯", "AI-Powered Matching", "Automatically connects brands with highest-performing publishers", ""),
         ("🛡️", "Fraud Detection", "Multi-layer protection against VPN masking, click farms, bot traffic", ""),
         ("⚡", "Real-Time Optimization", "Predictive campaign scoring, automated budget reallocation", "")]
for i, (icon, title, body, stat) in enumerate(cards):
    add_card(slide, Inches(0.8 + 4.1 * i), Inches(4.6), Inches(3.8), Inches(2.2),
             icon, title, body, stat)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 03 — TRACTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 3)
add_kicker(slide, "Proof — Traction", left=Inches(1.5), top=Inches(0.6))
add_title_centered(slide, "50+ Enterprise Brands.",
                   "$50M+ in Sales Driven.", top=Inches(1.0))

add_textbox(slide, Inches(2.5), Inches(2.5), Inches(8), Inches(0.5),
            "Real revenue, real clients, real results. Not a concept — a proven business growing 35% year-over-year.",
            font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

stats = [("870+", "Publishers", "LARGEST IN MENA"),
         ("$50M+", "Sales Driven", "FOR BRANDS"),
         ("8.2%", "Conversion Rate", "3–8× INDUSTRY"),
         ("87%", "Retention", "+20pts ABOVE AVG"),
         ("35%", "YoY Growth", "2.5× MARKET"),
         ("1.2M+", "Qualified Leads", "GENERATED")]
for i, (val, label, badge) in enumerate(stats):
    add_stat_box(slide, Inches(0.5 + 2.1 * i), Inches(3.2), val, label, badge)

# Brand logos as text list
brands_line1 = "noon · Talabat · NAMSHI · DAMAC · Trendyol · OSN · StarzPlay · HYCM · Evest"
brands_line2 = "eyewa · Homzmart · 6thStreet · OUNASS · SIVVI · Mumzworld · Level Shoes · NBA Store"
brands_line3 = "Calo · DailyMealz · MultiBank · hummel · Courir · Go Sport · Nice One · Temu · Styli"

txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.7), Inches(2))
tf = txBox.text_frame
tf.word_wrap = True
add_paragraph(tf, "Trusted by MENA's leading brands across 9 verticals", 9, TEXT_DIM,
              alignment=PP_ALIGN.CENTER, font_name="Consolas")
add_paragraph(tf, brands_line1, 12, TEXT_LIGHT, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
add_paragraph(tf, brands_line2, 12, TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
add_paragraph(tf, brands_line3, 12, TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 04 — PROBLEM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 4)
add_kicker(slide, "The Problem")
add_title(slide, "15% E-Commerce Growth.",
          "2% Affiliate Capture.")

add_subtitle(slide,
    "MENA's digital economy is booming but the performance marketing infrastructure is broken — or simply doesn't exist.")

# Two-column problem layout
brand_problems = [
    "❌ 80% of MENA brands have zero affiliate programs",
    "❌ CAC 335× higher than affiliate benchmarks",
    "❌ No reliable platform to find publishers",
    "❌ Fraud rates 3–5× global average",
    "❌ Manual tracking via spreadsheets & WhatsApp",
    "❌ No Arabic-language optimization tools",
]
pub_problems = [
    "❌ Zero reliable MENA-native platforms",
    "❌ Payment delays of 60–120 days",
    "❌ No data transparency or optimization tools",
    "❌ Western tools ignore Arabic content",
    "❌ No local payment rail support",
    "❌ Fragmented market — no single platform",
]

# For Brands column
add_textbox(slide, Inches(0.8), Inches(3.6), Inches(3), Inches(0.3),
            "🏢  FOR BRANDS", font_size=11, color=RED, bold=True, font_name="Consolas")
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(5.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
for prob in brand_problems:
    add_paragraph(tf, prob, 11, TEXT_LIGHT, space_before=Pt(3), space_after=Pt(3))

# For Publishers column
add_textbox(slide, Inches(7), Inches(3.6), Inches(3), Inches(0.3),
            "📱  FOR PUBLISHERS", font_size=11, color=RED, bold=True, font_name="Consolas")
txBox = slide.shapes.add_textbox(Inches(7), Inches(4.0), Inches(5.5), Inches(3))
tf = txBox.text_frame
tf.word_wrap = True
for prob in pub_problems:
    add_paragraph(tf, prob, 11, TEXT_LIGHT, space_before=Pt(3), space_after=Pt(3))

# Result bar
add_textbox(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.5),
            "THE RESULT: Billions in ad spend → Unmeasurable ROI → Brands overpay → Publishers underearned → Market stays underdeveloped",
            font_size=11, color=RED, alignment=PP_ALIGN.CENTER)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 05 — SOLUTION
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 5)
add_kicker(slide, "The Solution")
add_title(slide, "One Platform.",
          "Complete Affiliate Ecosystem.")

add_subtitle(slide,
    "Perff AI replaces spreadsheets, WhatsApp groups, and manual tracking with an AI-powered platform "
    "that automates the entire affiliate lifecycle.")

cards = [
    ("🎯", "AI Publisher Matching", "ML analyzes publisher audiences and performance for optimal pairings", "8.2% CVR"),
    ("🛡️", "Fraud Protection", "Multi-layer MENA-trained detection — VPN, click farms, bots", "3–5× Better"),
    ("⚡", "Smart Optimization", "Real-time budget reallocation, predictive scoring, auto-reporting", "7.0× ROAS"),
    ("📊", "Transparent Analytics", "Real-time dashboards, full attribution, conversion tracking", "Real-Time"),
    ("💰", "Reliable Payments", "On-time via bank transfer, PayPal, local methods. #1 pain solved", "64% Organic"),
    ("🌍", "MENA-Native", "Arabic-first, local currency, regional consumer intelligence", "4 Markets"),
]
for i, (icon, title, body, stat) in enumerate(cards):
    row = i // 3
    col = i % 3
    add_card(slide, Inches(0.8 + 4.1 * col), Inches(3.6 + 2.5 * row),
             Inches(3.8), Inches(2.2), icon, title, body, stat)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 06 — REVENUE MODEL
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 6)
add_kicker(slide, "Revenue Model")
add_title(slide, "How We",
          "Make Money.")

add_subtitle(slide,
    "Performance-based revenue with built-in alignment — we only make money when our clients get results.")

# Flow
flow_labels = [("🏢 Brand Pays", "Per conversion (CPA/CPS)", "$100"),
               ("🧠 Perff AI Keeps", "Performance margin", "$20–35"),
               ("📱 Publisher Earns", "Commission payout", "$65–80")]
for i, (label, sub, val) in enumerate(flow_labels):
    x = Inches(1.5 + 4 * i)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(3.4), Inches(3), Inches(1.2))
    shape.fill.solid()
    if i == 1:
        shape.fill.fore_color.rgb = GREEN
        lbl_color = DARK_BG
        val_color = DARK_BG
    else:
        shape.fill.fore_color.rgb = DARK2_BG
        lbl_color = WHITE
        val_color = GREEN
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.2), Inches(3.5), Inches(2.6), Inches(0.3),
                label, font_size=14, color=lbl_color, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(3.8), Inches(2.6), Inches(0.25),
                sub, font_size=9, color=lbl_color if i == 1 else TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), Inches(4.1), Inches(2.6), Inches(0.35),
                val, font_size=20, color=val_color, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

    if i < 2:
        add_textbox(slide, x + Inches(3.1), Inches(3.7), Inches(0.5), Inches(0.4),
                    "→", font_size=24, color=GREEN, alignment=PP_ALIGN.CENTER)

# Revenue stream cards
streams = [
    ("PRIMARY — LIVE", GREEN, "Performance Margin",
     "Brand pays per conversion. Perff AI keeps the spread. 20–35% gross margin on every transaction.", "90%+ of Revenue"),
    ("SCALING — BUILDING", ORANGE, "SaaS Features",
     "Premium AI tools — predictive scoring, fraud detection, API access. Subscription-based recurring revenue.", "Launching 2026"),
    ("HIGH-TOUCH — LIVE", BLUE, "Managed Services",
     "Full-service affiliate program management for enterprise brands entering MENA.", "Select Clients"),
]
for i, (kicker, kcolor, title, body, stat) in enumerate(streams):
    x = Inches(0.8 + 4.1 * i)
    y = Inches(5.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(3.8), Inches(2.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.15), y + Inches(0.15), Inches(3.5), Inches(0.2),
                kicker, font_size=8, color=kcolor, bold=True, font_name="Consolas")
    add_textbox(slide, x + Inches(0.15), y + Inches(0.4), Inches(3.5), Inches(0.25),
                title, font_size=13, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.7), Inches(3.5), Inches(0.7),
                body, font_size=9, color=TEXT_DIM)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.5), Inches(3.5), Inches(0.3),
                stat, font_size=18, color=GREEN, bold=True, font_name="Consolas")

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 07 — MARKET OPPORTUNITY
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 7)
add_kicker(slide, "Market Opportunity", left=Inches(1.5), top=Inches(0.6))
add_title_centered(slide, "A $30 Billion Market.",
                   "MENA Is the Fastest-Growing Piece.", top=Inches(1.0))

add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.5),
            "Global affiliate marketing is projected to reach $30B by 2031. MENA's share is growing 25%+ annually.",
            font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# TAM / SAM / SOM boxes
markets = [("TOTAL ADDRESSABLE MARKET", "$30B", "Global affiliate marketing by 2031", "7.7% CAGR"),
           ("SERVICEABLE ADDRESSABLE MARKET", "$1.5B", "MENA performance marketing", "25%+ YoY GROWTH"),
           ("SERVICEABLE OBTAINABLE MARKET", "$200M", "3-year target (UAE + KSA + Egypt)", "OUR BEACHHEAD")]
for i, (label, val, desc, badge) in enumerate(markets):
    x = Inches(0.8 + 4.1 * i)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(3.3), Inches(3.8), Inches(1.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    kcolor = GREEN if i == 2 else TEXT_DIM
    add_textbox(slide, x + Inches(0.15), Inches(3.4), Inches(3.5), Inches(0.25),
                label, font_size=8, color=kcolor, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_textbox(slide, x + Inches(0.15), Inches(3.7), Inches(3.5), Inches(0.6),
                val, font_size=36, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")
    add_textbox(slide, x + Inches(0.15), Inches(4.3), Inches(3.5), Inches(0.25),
                desc, font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), Inches(4.6), Inches(3.5), Inches(0.2),
                badge, font_size=8, color=GREEN, alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Bottom cards
market_cards = [("📈", "15% E-com Growth", "Fastest major region globally"),
                ("👥", "400M+ Under 30", "Digital-native consumers"),
                ("📱", ">90% Smartphones", "GCC penetration rate"),
                ("🏗️", "Vision 2030", "KSA e-commerce push")]
for i, (icon, title, body) in enumerate(market_cards):
    add_card(slide, Inches(0.8 + 3.1 * i), Inches(5.5), Inches(2.8), Inches(1.4),
             icon, title, body)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 08 — COMPETITIVE LANDSCAPE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 8)
add_kicker(slide, "Competitive Landscape")
add_title(slide, "No One Owns",
          "MENA Affiliate Marketing.")

add_subtitle(slide,
    "The only question is whether the winner will be a MENA-native platform or a Western player adapting late.")

headers = ["Platform", "MENA Focus", "AI/ML", "Arabic First", "Publisher Network", "Fraud Detection"]
rows = [
    ["Perff AI", "✅ 100%", "✅ Core", "✅ Yes", "✅ 870+ MENA", "✅ MENA-trained"],
    ["CJ Affiliate", "❌ Minimal", "⚠️ Basic", "❌ No", "⚠️ Global, few MENA", "⚠️ Global only"],
    ["Rakuten", "❌ None", "⚠️ Basic", "❌ No", "⚠️ Global focus", "⚠️ Not localized"],
    ["Impact.com", "⚠️ Entering", "✅ Strong", "❌ No", "⚠️ Growing", "⚠️ Global models"],
    ["Awin", "❌ None", "⚠️ Basic", "❌ No", "⚠️ EU/US-focused", "⚠️ Not MENA-tuned"],
    ["ArabyAds", "⚠️ Partial", "❌ None", "⚠️ Partial", "⚠️ Smaller", "❌ None"],
    ["Admitad", "⚠️ Some", "❌ Basic", "❌ No", "⚠️ Global", "⚠️ Basic"],
]
col_widths = [Inches(1.6), Inches(1.5), Inches(1.3), Inches(1.3), Inches(2.2), Inches(2.2)]
add_table(slide, Inches(0.6), Inches(3.5), Inches(12), col_widths, headers, rows)

# Moat bar
add_textbox(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
            "OUR MOAT: 870+ MENA publishers · MENA-trained AI · Arabic-first platform · "
            "50+ enterprise relationships · 250+ campaign dataset · Local payment infrastructure",
            font_size=10, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 09 — TECHNOLOGY & FLYWHEEL
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 9)
add_kicker(slide, "Technology & Flywheel", left=Inches(1.5), top=Inches(0.6))
add_title_centered(slide, "The AI Gets Smarter",
                   "With Every Campaign.", top=Inches(1.0))

add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.6),
            "Every campaign trains our models. Every publisher adds data. Every brand deepens our understanding. "
            "This creates a compounding advantage competitors would need years to replicate.",
            font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Flywheel nodes
fw_nodes = ["🏢 More Brands", "📊 More Data", "🧠 Smarter AI", "📈 Better Results", "📱 More Publishers"]
for i, label in enumerate(fw_nodes):
    x = Inches(0.5 + 2.5 * i)
    is_center = i == 2
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(3.5), Inches(2.2), Inches(0.9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN if is_center else DARK2_BG
    shape.line.fill.background()

    lbl_color = DARK_BG if is_center else WHITE
    add_textbox(slide, x + Inches(0.1), Inches(3.6), Inches(2), Inches(0.5),
                label, font_size=13, color=lbl_color, bold=True, alignment=PP_ALIGN.CENTER)

    if i < 4:
        add_textbox(slide, x + Inches(2.1), Inches(3.7), Inches(0.5), Inches(0.3),
                    "→", font_size=18, color=GREEN, alignment=PP_ALIGN.CENTER)

# Two column AI features
ai_live = [
    "✅ Fraud detection MVP (MENA-specific patterns)",
    "✅ Predictive campaign scoring",
    "✅ Automated performance reporting",
    "✅ Budget optimization recommendations",
    "✅ Publisher-brand matching algorithms",
]
ai_planned = [
    "🔨 Multi-layer fraud detection v2",
    "🔨 Adaptive real-time budget allocation",
    "🔨 Publisher recommendation engine",
    "🔨 Self-serve brand dashboard with AI assistant",
    "🔨 Open API for programmatic access",
    "🔨 NLP-based Arabic content analysis",
]

add_textbox(slide, Inches(0.8), Inches(4.7), Inches(4), Inches(0.3),
            "AI — LIVE TODAY", font_size=9, color=GREEN, bold=True, font_name="Consolas")
txBox = slide.shapes.add_textbox(Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.2))
tf = txBox.text_frame
tf.word_wrap = True
for item in ai_live:
    add_paragraph(tf, item, 11, TEXT_LIGHT, space_before=Pt(2), space_after=Pt(2))

add_textbox(slide, Inches(7), Inches(4.7), Inches(4), Inches(0.3),
            "AI — PLANNED (POST-INVESTMENT)", font_size=9, color=ORANGE, bold=True, font_name="Consolas")
txBox = slide.shapes.add_textbox(Inches(7), Inches(5.0), Inches(5.5), Inches(2.5))
tf = txBox.text_frame
tf.word_wrap = True
for item in ai_planned:
    add_paragraph(tf, item, 11, TEXT_LIGHT, space_before=Pt(2), space_after=Pt(2))

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10 — EXPANSION ROADMAP
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 10)
add_kicker(slide, "Expansion Roadmap")
add_title(slide, "UAE → KSA → Egypt →",
          "Pan-MENA.")

add_subtitle(slide,
    "Proven playbook in UAE. Now expanding to the region's two largest markets — Saudi Arabia and Egypt.")

roadmap = [
    ("2022", "Foundation — UAE Launch",
     "Founded in Abu Dhabi · First brand clients · 200+ publishers · ~$200K revenue"),
    ("2023", "Growth — Enterprise Penetration",
     "Signed noon, Talabat, DAMAC · Revenue 3× to $600K · 500+ publishers · Fraud detection MVP"),
    ("2024", "Scale — Market Leadership",
     "Revenue 2× to $1.2M · 50+ enterprise brands · 620+ publishers · 17.5% market share (UAE/KSA)"),
    ("2025", "Investment — Pre-Seed Round",
     "$1.8M projected revenue · Self-serve dashboard · Begin KSA expansion · Raise $8M pre-seed"),
    ("2026", "Expansion — KSA + Product",
     "Riyadh office · Full AI platform launch · 1,500+ publishers · $3.5M revenue · API launch"),
    ("2027", "Scale — Pan-MENA",
     "Egypt launch · 2,000+ publishers · $8–9M revenue · SaaS revenue stream live"),
]

y = Inches(3.5)
for year, title, desc in roadmap:
    add_textbox(slide, Inches(0.8), y, Inches(1), Inches(0.35),
                year, font_size=16, color=GREEN, bold=True, font_name="Consolas")
    add_textbox(slide, Inches(2), y, Inches(3), Inches(0.3),
                title, font_size=12, color=WHITE, bold=True)
    add_textbox(slide, Inches(5.2), y, Inches(7.5), Inches(0.35),
                desc, font_size=10, color=TEXT_DIM)
    y += Inches(0.55)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11 — FINANCIAL PERFORMANCE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 11)
add_kicker(slide, "Financial Performance", left=Inches(1.5), top=Inches(0.6))
add_title_centered(slide, "Revenue 6× in 3 Years.",
                   "Self-Funded.", top=Inches(1.0))

add_textbox(slide, Inches(2), Inches(2.4), Inches(9), Inches(0.5),
            "Built to $1.8M revenue with zero outside cash. Capital-efficient founders who know how to do more with less.",
            font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Bar chart (visual using rectangles)
rev_data = [("2022", "$200K", 0.12), ("2023", "$600K", 0.35),
            ("2024", "$1.2M", 0.6), ("2025E", "$1.8M", 0.75),
            ("2026P", "$3.5M", 0.88), ("2027P", "$8–9M", 1.0)]

max_bar_h = Inches(2.5)
bar_bottom = Inches(5.5)

for i, (year, val, pct) in enumerate(rev_data):
    x = Inches(1.5 + 1.8 * i)
    bar_h = int(max_bar_h * pct)
    is_projected = "E" in year or "P" in year

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   x, bar_bottom - bar_h, Inches(1.2), bar_h)
    shape.fill.solid()
    if is_projected:
        shape.fill.fore_color.rgb = RGBColor(0x2A, 0x5E, 0x4A)
    else:
        shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()

    # Value label above bar
    add_textbox(slide, x - Inches(0.1), bar_bottom - bar_h - Inches(0.3),
                Inches(1.4), Inches(0.3),
                val, font_size=12, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")
    # Year label below
    add_textbox(slide, x, bar_bottom + Inches(0.05), Inches(1.2), Inches(0.25),
                year, font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Bottom metric cards
metrics = [("Gross Margin", "20–35%", "Performance margin, growing with AI"),
           ("Net Revenue Retention", "~120%", "Existing clients increase spend YoY"),
           ("Revenue per Employee", "$200K", "$1.8M revenue ÷ 9 team members")]
for i, (title, val, desc) in enumerate(metrics):
    x = Inches(0.8 + 4.1 * i)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(6.0), Inches(3.8), Inches(1.1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.15), Inches(6.05), Inches(3.5), Inches(0.25),
                title, font_size=11, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.15), Inches(6.3), Inches(3.5), Inches(0.35),
                val, font_size=20, color=GREEN, bold=True, font_name="Consolas")
    add_textbox(slide, x + Inches(0.15), Inches(6.65), Inches(3.5), Inches(0.3),
                desc, font_size=9, color=TEXT_DIM)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12 — TEAM
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 12)
add_kicker(slide, "The Team", left=Inches(1.5), top=Inches(0.5))
add_title_centered(slide, "Built by MENA Insiders",
                   "Who Know This Market.", top=Inches(0.9))

add_textbox(slide, Inches(2), Inches(2.3), Inches(9), Inches(0.4),
            "10+ years of MENA performance marketing experience. A previous exit. Deep brand relationships.",
            font_size=13, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Founders
founders = [
    ("RH", "Rima Hani", "Co-Founder & CEO",
     "10+ years leading performance marketing across MENA. Built and scaled affiliate programs for the region's largest brands.",
     "Ex-Equiti · Ex-DigiZag · Abu Dhabi"),
    ("FA", "Fahad Alsaedi", "Co-Founder & CTO",
     "Serial technical founder with a previous exit. Former engineering lead at Nana (nana.sa). Deep expertise in scalable platforms & AI/ML.",
     "1× Exit · Ex-Nana · MBIT · Saudi Arabia"),
    ("SC", "Sadah Chanti", "Co-Founder & COO",
     "Bilingual Accounts Director & Digital Marketing Expert with 17+ years leading 20+ person teams. Blends strategic thinking, data, and AI. MPA + MIS.",
     "17+ Years · Accounts Director · Digital Marketing · Riyadh, KSA"),
]
for i, (initials, name, role, bio, tags) in enumerate(founders):
    x = Inches(0.5 + 4.2 * i)
    y = Inches(2.9)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, y, Inches(3.9), Inches(2.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    # Green top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.9), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = GREEN
    line.line.fill.background()

    # Avatar
    avatar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    x + Inches(0.2), y + Inches(0.2), Inches(0.6), Inches(0.6))
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = GREEN
    avatar.line.fill.background()
    add_textbox(slide, x + Inches(0.2), y + Inches(0.3), Inches(0.6), Inches(0.4),
                initials, font_size=18, color=DARK_BG, bold=True, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, x + Inches(0.95), y + Inches(0.2), Inches(2.8), Inches(0.3),
                name, font_size=16, color=WHITE, bold=True)
    add_textbox(slide, x + Inches(0.95), y + Inches(0.5), Inches(2.8), Inches(0.25),
                role, font_size=11, color=GREEN, bold=True)
    add_textbox(slide, x + Inches(0.2), y + Inches(0.9), Inches(3.5), Inches(0.9),
                bio, font_size=9, color=TEXT_LIGHT)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.9), Inches(3.5), Inches(0.5),
                tags, font_size=8, color=TEXT_DIM)

# Core team
add_textbox(slide, Inches(4.5), Inches(5.8), Inches(4), Inches(0.25),
            "CORE TEAM", font_size=9, color=TEXT_DIM, bold=True,
            alignment=PP_ALIGN.CENTER, font_name="Consolas")

team = [("Amer Kilani", "General Manager", "Amman"),
        ("Abeer Hani", "Head of Communications", "Riyadh"),
        ("Fairouz Hamdan", "Social Media Manager", "Amman"),
        ("Samer Baibers", "Sr. Affiliate Manager", "Amman"),
        ("Mohammed Alshaikh", "Sr. Affiliate Manager", "Gaza"),
        ("Lilian Alfar", "Sr. Affiliate Account Exec", "Amman")]
for i, (name, role, loc) in enumerate(team):
    x = Inches(0.5 + 2.15 * i)
    add_textbox(slide, x, Inches(6.15), Inches(2), Inches(0.2),
                name, font_size=10, color=WHITE, bold=True)
    add_textbox(slide, x, Inches(6.35), Inches(2), Inches(0.2),
                role, font_size=8, color=GREEN)
    add_textbox(slide, x, Inches(6.5), Inches(2), Inches(0.2),
                f"📍 {loc}", font_size=8, color=TEXT_DIM)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13 — THE ASK
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 13)
add_kicker(slide, "The Ask", left=Inches(1.5), top=Inches(0.6))
add_title_centered(slide, "$8,000,000 Pre-Seed",
                   "to Build the Category.", top=Inches(1.0))

add_textbox(slide, Inches(2), Inches(2.5), Inches(9), Inches(0.5),
            "We're raising $8M to build the definitive AI-powered affiliate marketing platform for MENA — "
            "before the window closes.",
            font_size=14, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)

# Use of Funds
funds = [
    (GREEN, "AI Development — ML Engineers, Infrastructure, Training", "40% · $3.2M"),
    (BLUE, "Regional Expansion — KSA (Riyadh), Egypt (Cairo), Sales Teams", "35% · $2.8M"),
    (ORANGE, "Marketing & Growth — Brand Acquisition, Publisher Growth, Events", "25% · $2.0M"),
]
for i, (color, label, val) in enumerate(funds):
    y = Inches(3.3 + 0.6 * i)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(3), y + Inches(0.05), Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_textbox(slide, Inches(3.3), y, Inches(5.5), Inches(0.3),
                label, font_size=12, color=TEXT_LIGHT)
    add_textbox(slide, Inches(9), y, Inches(2), Inches(0.3),
                val, font_size=12, color=GREEN, bold=True, font_name="Consolas")

# 12-Month / 24-Month milestone boxes
milestones_12 = [
    "✅ Full Perff AI platform GA launch",
    "✅ 3 ML engineers hired and shipping",
    "✅ KSA market entry — Riyadh office",
    "✅ 1,500+ publisher network",
    "✅ Self-serve brand dashboard live",
    "✅ $4M+ revenue run rate",
]
milestones_24 = [
    "✅ Egypt market entry — Cairo office",
    "✅ 2,000+ publisher network",
    "✅ Full AI stack deployed (fraud v2, scoring, API)",
    "✅ SaaS revenue stream live",
    "✅ $8–9M revenue run rate",
    "✅ Series A ready",
]

for col, (title, items) in enumerate(
        [("12-MONTH MILESTONES", milestones_12),
         ("24-MONTH MILESTONES", milestones_24)]):
    x = Inches(0.8 + 6.3 * col)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(5.0), Inches(5.8), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.2), Inches(5.1), Inches(5), Inches(0.25),
                title, font_size=9, color=GREEN, bold=True, font_name="Consolas")
    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(5.4), Inches(5.4), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for item in items:
        add_paragraph(tf, item, 10, TEXT_LIGHT, space_before=Pt(2), space_after=Pt(2))

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14 — WHY NOW
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 14)
add_kicker(slide, "Why Now")
add_title(slide, "The Window Is",
          "18–24 Months.")

add_subtitle(slide,
    "Five macro forces are converging to create a once-in-a-decade opportunity in MENA affiliate marketing.")

why_now = [
    ("🛒", "E-Commerce Boom", "MENA e-commerce growing 15% YoY. $50B+ market. Every brand needs efficient acquisition."),
    ("🏗️", "Vision 2030", "Saudi Arabia investing billions in digital infrastructure. New e-commerce brands launching monthly."),
    ("🧠", "AI Inflection", "Gen-AI and ML tools now accessible. Our 250+ campaign dataset gives us a head start."),
    ("💸", "CAC Crisis", "Facebook/Google CPAs rising 30%+ YoY in MENA. Brands seeking performance alternatives."),
    ("🌍", "No Incumbent", "No dominant affiliate platform in MENA. Western players haven't localized."),
    ("⏰", "First-Mover Window", "International platforms will enter MENA within 18–24 months. We must lead now."),
]
for i, (icon, title, body) in enumerate(why_now):
    row = i // 3
    col = i % 3
    add_card(slide, Inches(0.8 + 4.1 * col), Inches(3.5 + 2.4 * row),
             Inches(3.8), Inches(2.1), icon, title, body)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 15 — CLOSING
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_slide_number(slide, 15)
add_kicker(slide, "Let's Build the Future of MENA Performance Marketing",
           left=Inches(1.5), top=Inches(1.0))

txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.3), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
run1 = p.add_run()
run1.text = "Turn MENA's Ad Spend Into\n"
run1.font.size = Pt(48)
run1.font.color.rgb = WHITE
run1.font.bold = True
run1.font.name = "Calibri"
run2 = p.add_run()
run2.text = "Measurable Growth."
run2.font.size = Pt(48)
run2.font.color.rgb = GREEN
run2.font.bold = True
run2.font.name = "Calibri"

# Key stats
closing_stats = [("Raising", "$8M"), ("Revenue", "$1.8M"), ("Brands", "50+"), ("Publishers", "870+")]
for i, (label, val) in enumerate(closing_stats):
    x = Inches(2.5 + 2.2 * i)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x, Inches(3.5), Inches(1.8), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK2_BG
    shape.line.fill.background()

    add_textbox(slide, x + Inches(0.1), Inches(3.55), Inches(1.6), Inches(0.2),
                label, font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), Inches(3.8), Inches(1.6), Inches(0.5),
                val, font_size=22, color=GREEN, bold=True,
                alignment=PP_ALIGN.CENTER, font_name="Consolas")

# Contact card
contact_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(4), Inches(4.8), Inches(5), Inches(1.8))
contact_shape.fill.solid()
contact_shape.fill.fore_color.rgb = DARK2_BG
contact_shape.line.fill.background()

add_textbox(slide, Inches(4.2), Inches(4.9), Inches(4.6), Inches(0.35),
            "Rima Hani", font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.2), Inches(5.25), Inches(4.6), Inches(0.25),
            "Co-Founder & CEO", font_size=11, color=GREEN, alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(4.2), Inches(5.5), Inches(4.6), Inches(0.25),
            "rima.hani@globalnetwork.global", font_size=11, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

contact_details = ("perff.net\n"
                   "USA (HQ) · Abu Dhabi, UAE · Amman, Jordan · Riyadh, KSA · Gaza, Palestine\n"
                   "UAE: +971 506856839 · Jordan: +962 796314887")
add_textbox(slide, Inches(4.2), Inches(5.8), Inches(4.6), Inches(0.8),
            contact_details, font_size=9, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Copyright
add_textbox(slide, Inches(3.5), Inches(6.8), Inches(6), Inches(0.3),
            "© 2026 Perff AI (Perff Inc.). Confidential & Proprietary.",
            font_size=10, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

add_confidential(slide)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
output_path = os.path.join(os.path.dirname(__file__),
                           "data-room", "02-pitch-deck",
                           "Perff-AI-Investor-Pitch-Deck.pptx")
prs.save(output_path)
print(f"✅ Pitch deck saved to: {output_path}")
print(f"   {len(prs.slides)} slides generated")
print(f"   Slide size: 13.33\" × 7.5\" (widescreen 16:9)")
